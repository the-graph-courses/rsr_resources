#!/usr/bin/env python3
"""Build a single-slide 16:9 PowerPoint with the cupcake t-test SVG embedded as a
vector picture (PNG fallback + svgBlip), so it stays crisp and 'Convert to Shape'-able."""
import re, shutil, zipfile, os
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
SVG = os.path.join(HERE, "cupcake_t_test_slide.svg")
PNG = os.path.join(HERE, "slide_fallback.png")
TMP = os.path.join(HERE, "_base.pptx")
OUT = os.path.join(HERE, "cupcake_t_test_slide.pptx")

# 16:9 slide = 13.333in x 7.5in
SW, SH = Emu(12192000), Emu(6858000)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide.shapes.add_picture(PNG, 0, 0, width=SW, height=SH)
prs.save(TMP)

# ---- inject the SVG with svgBlip fallback ----
with open(SVG, "rb") as f:
    svg_bytes = f.read()

zin = zipfile.ZipFile(TMP, "r")
names = zin.namelist()
items = {n: zin.read(n) for n in names}
zin.close()

# 1) content types: ensure svg default
ct = items["[Content_Types].xml"].decode("utf-8")
if "image/svg+xml" not in ct:
    ct = ct.replace("</Types>",
        '<Default Extension="svg" ContentType="image/svg+xml"/></Types>')
items["[Content_Types].xml"] = ct.encode("utf-8")

# 2) add svg media
items["ppt/media/image2.svg"] = svg_bytes

# 3) slide rels: add relationship to the svg
rels_name = "ppt/slides/_rels/slide1.xml.rels"
rels = items[rels_name].decode("utf-8")
SVG_RID = "rId100"
rel = ('<Relationship Id="%s" '
       'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
       'Target="../media/image2.svg"/>' % SVG_RID)
rels = rels.replace("</Relationships>", rel + "</Relationships>")
items[rels_name] = rels.encode("utf-8")

# 4) slide xml: add svgBlip extension inside the existing <a:blip>
slide_name = "ppt/slides/slide1.xml"
xml = items[slide_name].decode("utf-8")
m = re.search(r'<a:blip\b[^>]*?r:embed="([^"]+)"[^>]*?/>', xml)
if not m:
    raise SystemExit("could not find self-closing a:blip in slide xml:\n" + xml[:2000])
png_rid = m.group(1)
ext = ('<a:extLst><a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">'
       '<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
       'r:embed="%s"/></a:ext></a:extLst>' % SVG_RID)
new_blip = '<a:blip r:embed="%s">%s</a:blip>' % (png_rid, ext)
xml = xml[:m.start()] + new_blip + xml[m.end():]
items[slide_name] = xml.encode("utf-8")

# 5) write final
with zipfile.ZipFile(OUT, "w", zipfile.ZIP_DEFLATED) as z:
    for n, data in items.items():
        z.writestr(n, data)

os.remove(TMP)
print("wrote", OUT, "png rId=", png_rid)
